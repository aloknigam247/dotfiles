"""
pptlib - reusable python-pptx helpers for themed 16:9 presentations.

Usage:
    from pptlib import Deck
    d = Deck()                      # default navy+teal theme
    s = d.slide()
    d.header(s, "Architecture", "System Overview", 3)
    d.bullet(d.textbox(s, 0.9, 2.0, 5, 4), "First point", first=True)
    d.img_fit(s, "img/arch.png", 6.3, 2.0, 6.3, 4.5)
    d.save(r"C:\\path\\Deck.pptx")

Override the theme by passing a dict of hex colors / font names:
    Deck(theme={"accent": "E11D48", "navy": "1A1A2E", "font": "Calibri"})
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from PIL import Image

ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}

DEFAULT_THEME = {
    "navy": "0F2A43", "navy_dk": "0A1F33", "accent": "14B8A6", "accent2": "2563EB",
    "light": "F5F7FA", "white": "FFFFFF", "text": "1E293B", "muted": "64748B",
    "warn": "D97706", "panel": "ECF1F7", "card_border": "E2E8F0",
    "subtle": "C8D6E5", "font": "Segoe UI", "mono": "Consolas",
}


def hex2rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class Deck:
    """A 16:9 presentation builder with a consistent theme."""

    def __init__(self, theme=None):
        self.t = {**DEFAULT_THEME, **(theme or {})}
        for k, v in self.t.items():
            if k not in ("font", "mono"):
                setattr(self, k.upper(), hex2rgb(v))
        self.FONT = self.t["font"]
        self.MONO = self.t["mono"]
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self._blank = self.prs.slide_layouts[6]

    # -- slide / save -----------------------------------------------------
    def slide(self):
        return self.prs.slides.add_slide(self._blank)

    def save(self, path):
        self.prs.save(path)
        return path

    # -- shapes -----------------------------------------------------------
    def rect(self, s, x, y, w, h, fill=None, line=None, line_w=1.0,
             shape=MSO_SHAPE.RECTANGLE, shadow=False):
        sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        if fill is None:
            sp.fill.background()
        else:
            sp.fill.solid()
            sp.fill.fore_color.rgb = fill
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line
            sp.line.width = Pt(line_w)
        # IMPORTANT: shadow.inherit=False inserts an empty <a:effectLst/>. Adding a
        # SECOND effectLst produces schema-invalid XML that PowerPoint refuses to
        # open (python-pptx is more lenient). Reuse the existing element instead.
        sp.shadow.inherit = False
        if shadow:
            spPr = sp._element.spPr
            ef = spPr.find(qn("a:effectLst"))
            if ef is None:
                ef = spPr.makeelement(qn("a:effectLst"), {})
                spPr.append(ef)
            sh = ef.makeelement(qn("a:outerShdw"),
                                {"blurRad": "90000", "dist": "38100",
                                 "dir": "5400000", "rotWithShape": "0"})
            clr = sh.makeelement(qn("a:srgbClr"), {"val": self.t["navy"]})
            clr.append(clr.makeelement(qn("a:alpha"), {"val": "28000"}))
            sh.append(clr)
            ef.append(sh)
        return sp

    def connector(self, s, x1, y1, x2, y2, color=None, width=1.5,
                  kind=MSO_CONNECTOR.STRAIGHT):
        cn = s.shapes.add_connector(kind, Inches(x1), Inches(y1),
                                    Inches(x2), Inches(y2))
        cn.line.color.rgb = color or self.ACCENT2
        cn.line.width = Pt(width)
        cn.shadow.inherit = False
        return cn

    # -- text -------------------------------------------------------------
    def textbox(self, s, x, y, w, h, anchor="top"):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = ANCHOR[anchor]
        for m in ("left", "right", "top", "bottom"):
            setattr(tf, f"margin_{m}", 0)
        return tf

    def _para(self, tf, first):
        if first and not tf.paragraphs[0].runs:
            return tf.paragraphs[0]
        return tf.add_paragraph()

    def run(self, p, text, size=18, color=None, bold=False, italic=False,
            mono=False, font=None):
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color if color is not None else self.TEXT
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font or (self.MONO if mono else self.FONT)
        return r

    def para(self, tf, first=False, align=None, space_before=0, space_after=0):
        p = self._para(tf, first)
        if align:
            p.alignment = ALIGN[align]
        p.space_before = Pt(space_before)
        p.space_after = Pt(space_after)
        return p

    def bullet(self, tf, text, size=18, color=None, bold=False, level=0,
               glyph="\u2022", gcolor=None, space_after=10, first=False):
        p = self.para(tf, first, space_after=space_after)
        p.level = level
        if glyph:
            self.run(p, glyph + "  ", size=size,
                     color=gcolor if gcolor is not None else self.ACCENT, bold=True)
        self.run(p, text, size=size, color=color, bold=bold)
        return p

    # -- images -----------------------------------------------------------
    def img_fit(self, s, path, bx, by, bw, bh, align="center", valign="middle"):
        """Insert an image scaled to fit inside the box, preserving aspect ratio."""
        iw, ih = Image.open(path).size
        ar, bar = iw / ih, bw / bh
        if ar > bar:
            w, h = bw, bw / ar
        else:
            w, h = bh * ar, bh
        x = {"left": bx, "center": bx + (bw - w) / 2, "right": bx + bw - w}[align]
        y = {"top": by, "middle": by + (bh - h) / 2, "bottom": by + bh - h}[valign]
        s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
        return x, y, w, h

    # -- composite layout -------------------------------------------------
    def header(self, s, kicker, title, num, footer_label="Presentation"):
        """Standard content-slide chrome: bg, accent bar, kicker, title, footer."""
        self.rect(s, 0, 0, 13.333, 7.5, fill=self.LIGHT)
        self.rect(s, 0, 0, 0.22, 7.5, fill=self.ACCENT)
        tf = self.textbox(s, 0.7, 0.42, 11.5, 0.4)
        self.run(self.para(tf, True), kicker.upper(), size=12, color=self.ACCENT, bold=True)
        tf2 = self.textbox(s, 0.7, 0.72, 11.9, 0.9)
        self.run(self.para(tf2, True), title, size=30, color=self.NAVY, bold=True)
        self.rect(s, 0.72, 1.55, 0.9, 0.06, fill=self.ACCENT2)
        self.footer(s, num, footer_label)

    def footer(self, s, num, label="Presentation"):
        tf = self.textbox(s, 0.7, 7.02, 8, 0.35)
        self.run(self.para(tf, True), label, size=10, color=self.MUTED)
        tf2 = self.textbox(s, 11.8, 7.02, 0.8, 0.35)
        self.run(self.para(tf2, True, align="right"), f"{num:02d}", size=10,
                 color=self.MUTED, bold=True)

    def node(self, s, x, y, w, h, text, fill=None, text_color=None, size=14,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True):
        """A labelled box for hand-built (native) diagrams. Returns the shape."""
        sp = self.rect(s, x, y, w, h, fill=fill or self.PANEL,
                       line=self.ACCENT, line_w=1.5, shape=shape, shadow=shadow)
        tf = sp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        self.run(self.para(tf, True, align="center"), text, size=size,
                 color=text_color or self.NAVY, bold=True)
        return sp
